from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE,XL_DATA_LABEL_POSITION, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT,MSO_ANCHOR, MSO_AUTO_SIZE
import pymssql
import pandas as pd
import os
import datetime
import dateutil
import re
from sqlalchemy import create_engine


def get_datelist(today):
    yst = datetime.date(today.year,1,1)
    if today - yst < datetime.timedelta(31):
        yst = datetime.date(today.year-1,1,1)
    datearray = pd.date_range(yst,today,freq='M')
    return datearray.astype(str).tolist()


# 画一个橙色虚线框，放置文字内容
def myRen(shapes,paragraph_str=[], left=Inches(8.5), top=Inches(2), width=Inches(4.5), height=Inches(5)):  # 创建一个橙色虚线框，未解决居中问题
    # left = top = width = height = Inches(1.0)
    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    # 背景色透明
    fill = shape.fill
    fill.background()
    # 线条为虚线
    line = shape.line
    line.color.rgb = RGBColor(210, 71, 38)
    line.dash_style = MSO_LINE.DASH
    line.width = Pt(1.5)
    # 更改形状大小
    # shape.height = Inches(4)
    # shape.width = Inches(2)
    text_frame = shape.text_frame
    text_frame.margin_bottom = Inches(0.3)
    text_frame.margin_left = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    if paragraph_str:
        for line in paragraph_str:
            p = text_frame.add_paragraph()
            p.line_spacing = 1.3
            run = p.add_run()
            run.text = line
            run.font.bold = True
            run.font.color.rgb = black
            run.font.size = Pt(20)
    return shape


# 竖向柱形图
def drawbar2(shapes, df, title, left=Inches(0.5), top=Inches(1), width=Inches(7), height=Inches(6.5), sorted=True, fontsize=18, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    # values = df.values.tolist()
    columns = df.columns.values.tolist()  # 列名
    df = df.dropna()
    if sorted:  # 默认对数值排序
        sort_item = columns[1]
        df = df.sort_values(by=sort_item)
    df_labels = df.iloc[:,0]  # 部门
    df_values = df.iloc[:,1:]  # 数值
    labels = df_labels.values.tolist()
    chart_data = ChartData()
    for i in range(df_values.shape[1]):
        values = df_values.iloc[:,i].values.tolist()
        column = columns[i+1]
        chart_data.add_series(column, values)
    chart_data.categories = labels
    # x, y, cx, cy = Inches(0.5), Inches(1), Inches(6), Inches(6.5)
    graphic_frame = shapes.add_chart(chart_type, left, top, width, height, chart_data)
    graphic_frame.chart.has_legend = True
    graphic_frame.chart.legend.position = XL_LEGEND_POSITION.CORNER  # 图例在右上角
    graphic_frame.chart.has_title = True
    graphic_frame.chart.chart_title.text_frame.clear()
    # graphic_frame.chart.has_table = True
    new_title = graphic_frame.chart.chart_title.text_frame.add_paragraph()
    new_title.text = title
    new_title.font.size = Pt(24)
    new_title.font.bold = True
    series = graphic_frame.chart.series[0]
    series.invert_if_negative = False  # 没什么卵用
    plot = graphic_frame.chart.plots[0]  # 取图表中第一个plot
    plot.has_data_labels = True  # 是否显示数据标签
    data_labels = plot.data_labels  # 数据标签控制类
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END  # 字体位置
    data_labels.font.bold = True
    if True in (df_values[df_values<1].count()/df_values.count()>0.9).tolist():
        data_labels.number_format = '0.0%'
    data_labels.font.size = Pt(fontsize)
    category_axis = graphic_frame.chart.category_axis  # 纵轴标签控制类
    category_axis.tick_labels.font.bold=True
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis.has_major_gridlines = False
    value_axis = graphic_frame.chart.value_axis  # 横轴值坐标标签控制类
    # value_axis.tick_labels.number_format = '0%'
    value_axis.has_minor_gridlines = False
    value_axis.has_major_gridlines = False
    return graphic_frame


# 横向柱形图
def drawbar(shapes, df, title, left=Inches(0.5), top=Inches(1), width=Inches(6), height=Inches(6.5), sorted=True, fontsize=18):
    # values = df.values.tolist()
    columns = df.columns.values.tolist()  # 列名
    df = df.dropna()
    if sorted:  # 默认对数值排序
        sort_item = columns[1]
        df = df.sort_values(by=sort_item)
    df_labels = df.iloc[:,0]  # 部门

    df_values = df.iloc[:,1:]  # 数值
    labels = df_labels.values.tolist()
    chart_data = ChartData()
    for i in range(df_values.shape[1]):
        values = df_values.iloc[:,i].values.tolist()
        column = columns[i+1]
        chart_data.add_series(column, values)
    chart_data.categories = labels
    # x, y, cx, cy = Inches(0.5), Inches(1), Inches(6), Inches(6.5)
    graphic_frame = shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data)
    graphic_frame.chart.has_title = True
    graphic_frame.chart.chart_title.text_frame.clear()
    new_title = graphic_frame.chart.chart_title.text_frame.add_paragraph()
    new_title.text = title
    new_title.font.size = Pt(24)
    new_title.font.bold = True
    series = graphic_frame.chart.series[0]
    series.invert_if_negative = False  # 没什么卵用
    for i in range(df_values.shape[0]):
        point = series.points[i]
        fill = point.format.fill
        fill.patterned()  # 此处不可用solid()，否则负值会被取反色
        fill.fore_color.rgb = RGBColor(210, 71, 38)  # orange
        fill.back_color.rgb = RGBColor(210, 71, 38)  # 背景色也被设置为橙色，背景色就是负值颜色
    plot = graphic_frame.chart.plots[0]  # 取图表中第一个plot
    plot.has_data_labels = True  # 是否显示数据标签
    data_labels = plot.data_labels  # 数据标签控制类
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END  # 字体位置
    data_labels.font.bold = True
    data_labels.number_format = '0.0%'
    data_labels.font.size = Pt(fontsize)
    category_axis = graphic_frame.chart.category_axis  # 纵轴标签控制类
    category_axis.tick_labels.font.bold=True
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    category_axis.has_major_gridlines = False
    value_axis = graphic_frame.chart.value_axis  # 横轴值坐标标签控制类
    value_axis.tick_labels.number_format = '0%'
    value_axis.has_minor_gridlines = False
    value_axis.has_major_gridlines = False
    return graphic_frame


# 画表格，百分比内容可以标记红色
def drawtable(shapes, df, title='', left=Inches(7.0), top=Inches(1.5), width=Inches(6.5), height=Inches(5.5), red_mark='below_sum'):
    '''red_mark: 可选:below_zero,below_sum.前者将负值标记为红色，后者将带百分号又低于合计值的标记为红色'''
    df = df.fillna('')
    rows = df.shape[0]+df.columns.nlevels
    cols = df.shape[1]
    values = df.values.tolist()
    values.insert(0,df.columns.values.tolist())  # 插入标题行
    sumline = values[-1]  # 取最后一行为合计值
    table = shapes.add_table(rows, cols, left, top, width, height).table
    # 填充单元格
    column_range = [i for i in range(0,df.columns.nlevels)]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r,c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # cell.text = values[r][c]
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.font.bold = True
            p.text = values[r][c]
            p.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            # p.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r in column_range:  # 标题行保持主题颜色
                pass
            else:  # 内容行改色
                color2 = RGBColor(248,215,205)  # 深
                color1 = RGBColor(252,236,232)  # 浅
                color = color1 if r % 2 == 0 else color2
                cell = table.cell(r,c)
                cell_value = values[r][c]
                cell.text = cell_value
                cell.fill.solid()
                cell.fill.fore_color.rgb = color
                if red_mark == 'below_zero':  # 低于零的值标记为红色
                    if '-' in cell_value and '%' in cell_value:
                        cell.text_frame.paragraphs[0].font.color.rgb = red
                elif red_mark == 'below_sum':  # 低于合计值的值标记为红色
                    if '%' in cell_value:
                        try:
                            if float(cell_value.replace('%','')) < float(sumline[c].replace('%','')):
                                cell.text_frame.paragraphs[0].font.color.rgb = red
                        except:
                            pass
    # 标题文本框
    tb = shapes.add_textbox(left, top-Inches(0.5), width, Inches(0.5))
    tb_tf = tb.text_frame
    tb_tf.clear()
    p = tb_tf.paragraphs[0]
    p.text = title
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.bold = True
    p.font.size = Pt(24)
    # new_title = tb_tf.add_paragraph()
    # new_title.text = title
    # new_title.font.size = Pt(24)
    # new_title.font.bold = True
    # new_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    return table


def wash_values(df):
    if '排名' in df.columns:
        df.loc[df[df['部门'] != '分公司'].index, '排名'] = df.loc[df[df['部门'] != '分公司'].index, '排名'].rank().astype(int)
        df.loc[12, '排名'] = '*'
    for column in df.columns:
        try:
            df[column] = (df[column] / 10000).astype(float).round(2)
        except:
            try:
                df[column] = df[column].apply(
                    lambda x: float(x.replace('%', '')) * 0.01)
            except:
                pass
    return df


def div_10000(df2):
    df = df2.copy()
    if '排名' in df.columns:
        df.loc[df[df['部门'] != '分公司'].index, '排名'] = df.loc[df[df['部门'] != '分公司'].index, '排名'].rank().astype(int)
        df.loc[12, '排名'] = '*'
    for column in df.columns:
        try:
            df[column] = (df[column] / 10000).astype(float).round(2)
        except:
            pass
    return df


def pct_2_float(df):  # 暴力转float，返回拷贝版本
    df2 = df.copy()
    for column in df2.columns:
        try:
            df2[column] = df2[column].map(
                lambda x: float(x.replace('%', '')) * 0.01,na_action='ignore')
        except:
            pass
    return df2


# 新建幻灯片，母版2：首页，母版3：标题页，母版0：带左右标题的正文页
def getshape(title_left='',title_right='',slide_index=0):
    slide = prs.slides.add_slide(prs.slide_layouts[slide_index])
    ph = slide.placeholders
    ph[0].text = title_left
    ph[1].text = title_right
    shapes = slide.shapes
    return shapes


# 左横柱图，右表格板式页面的快捷方式
def slide_with_bar_and_table(df,table_title,bar_title,table_column,bar_column,title_left,title_right=''):
    df_table = div_10000(df)[table_column].astype(str)  # 转万元
    df_bar = pct_2_float(df_table).loc[df[~df['部门'].isin(['分公司','合计'])].index,bar_column]
    shapes = getshape(title_left,title_right)
    drawtable(shapes, df_table, table_title)
    drawbar(shapes, df_bar, bar_title)


# 约束过长百分比的小数位数，限制单元格大小
def format_pct(item):
    try:
        value = float(str(item).replace('/','').replace('%',''))*0.01
        if len(str(item))<=6:
            out = format(value,'0.1%')
        else:
            out = format(value,'0.0%')
    except:
        out = item
    return out


# 获得最新一期市场报表，返回datetime格式日期及文件名
def newest_market_report():
    path = os.path.join(os.getcwd(),r'market_report\输出报表')
    fnlt = os.listdir(path)
    pattern = re.compile(r'(\d{4}年\d{1,2}月).*xlsx')
    datedict = {}
    for item in fnlt:
        sch = re.search(pattern,item)
        if sch:
            d_str = sch.group(1)
            d = datetime.datetime.strptime(d_str,'%Y年%m月')
            datedict.update({d:item})
    datemax = max(datedict.keys())
    fnmax = os.path.join(path,datedict[datemax])
    return datemax, fnmax


def main():
    outfn = 'test.pptx'
    # 首页
    getshape('2018', '09', 2)
    # 标题页1
    getshape('01', '车险经营概况', 3)
    # 保费增速，当年及当月
    reportfn = datelist[-1]+'报表查询系统车险增速.xlsx'
    if not os.path.exists(reportfn):
        os.system('python reportGD_rate.py')
    reportGD = pd.read_excel(reportfn, header=[0,1,2])
    # 车均
    with open('车均算法.sql','r') as fn:
        sql = fn.read()
    chejun = pd.read_sql(sql,engine)
    chejun.set_index('index',inplace=True)
    chejun = chejun.rename_axis({'thismonth':'当月','thisyear':'当年'})
    chejun['同比'] = chejun['今年车均']/chejun['去年车均']-1
    for timeflag in ['当月','当年']:  # 当年、当月两张幻灯片
        shapes = getshape('车险经营概况', '整体保费收入')
        thisyear = reportGD[timeflag].loc[datelist[-1], (slice(None),('今年','去年'))].unstack().reset_index()
        bar_title = year_str+allmonth_str+'保费收入（万）' if timeflag == '当年' else year_str+month_str+'保费收入（万）'  # 柱图标题
        drawbar2(shapes, thisyear, bar_title,
                fontsize=14)
        thisyear = thisyear.set_index('index')
        thisyear = thisyear.astype(float).round(2)
        thisyear['同比'] = thisyear['今年']/thisyear['去年']-1
        ss, sf, ss_rate, sf_rate = thisyear.loc['实收保费','今年'], thisyear.loc['收付费保费','今年'], thisyear.loc['实收保费','同比']\
            , thisyear.loc['收付费保费','同比']
        cj, cj_rate = chejun.loc[timeflag,'今年车均'], chejun.loc[timeflag,'同比']
        plus_minus1 = lambda x: '上升' if x >= 0 else '下降'  # 判断同比正负
        pct = lambda x:format(x,'0.2%') if isinstance(x,float) else x  # float转为百分比
        text = '实收保费：;{ss}万，同比{ss_rate} ;收付费保费：;{sf}，同比{sf_rate} ;车均保费：;{cj}元，同比{cj_rate} ;'\
            .format(ss=ss,sf=sf
            ,ss_rate=plus_minus1(ss_rate)+pct(abs(ss_rate))
            ,sf_rate=plus_minus1(sf_rate)+pct(abs(sf_rate))
            ,cj=cj,cj_rate=plus_minus1(cj_rate)+pct(abs(cj_rate)))
        shape = myRen(shapes,text.split(';'))
    # 行业报表
    datemax, fnmax = newest_market_report()
    year_str_market = str(datemax.year) + '年'
    month_str_market = str(datemax.month) + '月'
    allmonth_str_market = '1-%s' % month_str_market if month_str_market != '1月' else month_str_market
    ex = pd.ExcelFile(fnmax)
    el = []
    for st in ex.sheet_names:
        s = ex.parse(st)
        line = s[s.保费收入.isin(['人保财险', '平安财险', '太保财险', '国寿财', '行业合计'])][['保费收入', '增速']]
        line.rename_axis({'增速': st + '增速'}, 1, inplace=True)
        el.append(line)
    first = el[0]
    for i in range(1, len(el)):
        first = pd.merge(first, el[i])
    f = lambda x:float(str(x).replace('%', '')) * 0.01 if '%' in str(x) else x  # 百分比转小数
    first = first.applymap(f)
    shapes = getshape('车险经营概况', '保费与对标—市场份额')
    drawbar2(shapes, first,year_str_market+month_str_market+'各主体车险保费增速',sorted=False)
    first = first.set_index('保费收入')
    first = first.T
    first['与行业差距'] = first['人保财险']-first['行业合计']
    first['与平安差距'] = first['人保财险'] - first['平安财险']
    first = first.T
    first = first.applymap(pct)  # 统一转换为百分比
    sc_month, rb_month, pa_month, div_sc_month, div_pa_month = first.loc['行业合计','当月增速'], first.loc['人保财险','当月增速'], first.loc['平安财险','当月增速'], first.loc['与行业差距','当月增速'], first.loc['与平安差距','当月增速']
    sc_year, rb_year, pa_year, div_sc_year, div_pa_year = first.loc['行业合计', '当年增速'], first.loc['人保财险', '当年增速'], \
                                                          first.loc['平安财险', '当年增速'], first.loc['与行业差距', '当年增速'], \
                                                          first.loc['与平安差距', '当年增速']
    plus_minus2 = lambda x:'领先' if '-' not in x else '落后'  # 根据百分比判断
    str_abs = lambda x:x.replace('-','') if isinstance(x,str) else x  # 去除负号
    text = '''{m}珠海市车险行业增速{sc_month}，我司增速{rb_month}，{p_or_m1}市场平均增速{div_sc_month}，{p_or_m2}平安{div_pa_month};
{allm}珠海市车险行业增速为{sc_year}，我司增速{rb_year}，{p_or_m3}市场{div_sc_year}，{p_or_m4}平安{div_pa_year}'''\
        .format(sc_month=sc_month,rb_month=rb_month,div_sc_month=str_abs(div_sc_month), div_pa_month=str_abs(div_pa_month)
                ,sc_year=sc_year,rb_year=rb_year,div_sc_year=str_abs(div_sc_year), div_pa_year=str_abs(div_pa_year)
                ,p_or_m1=plus_minus2(div_sc_month),p_or_m2=plus_minus2(div_pa_month),p_or_m3=plus_minus2(div_sc_year)
                ,p_or_m4=plus_minus2(div_pa_year)
                ,m=month_str_market,allm=allmonth_str_market)
    shape = myRen(shapes,text.split(';'))
    # 日通报内容
    # 标题页2
    getshape('02', '车险业务点评', 3)
    # 运行日通报
    sql = """exec car.dbo.dailyreport_ppt '%s'""" % deskdate
    # sql = """exec car.dbo.dailyreport_ppt '%s'""" % '20180105'
    conn = pymssql.connect(host='56.49.33.246', port='1433', user='isp', password='pIcc4404', database='car')
    cursor = conn.cursor(as_dict=True)
    cursor.execute(sql)
    titles = ['保费增速','家用车保费增速','当月新续转增速','当年新续转增速','续保率','家用车当月续保率','提前签单占比','当月渠道增速','当年渠道增速']
    title_left = '车险经营概况'
    # 保费增速
    query = cursor.fetchall()
    # 保费报表
    df = pd.DataFrame(query)
    # slide_with_bar_and_table(df,'当月部门起保增速明细', '当月部门起保增速',['排名','部门','当月起保保费','当月起保增速'],['部门','当月起保增速'],title_left,'当月起保增速')
    # slide_with_bar_and_table(df, '当年部门起保增速明细', '当年部门起保增速', ['排名', '部门', '当年起保保费', '当年起保保费增速'], ['部门', '当年起保保费增速'],title_left,'当年起保增速')
    # slide_with_bar_and_table(df, '当月部门签单增速明细', '当月部门签单增速', ['排名', '部门', '当月签单保费', '当月签单增速'], ['部门', '当月签单增速'],title_left,'当月签单增速')
    # slide_with_bar_and_table(df, '当年部门签单增速明细', '当年部门签单增速', ['排名', '部门', '当年签单保费', '当年签单保费增速'], ['部门', '当年签单保费增速'],title_left,'当年签单增速')
    # 家用增速跳过
    df = pd.DataFrame(cursor.fetchall())
    # 当月新续转增速
    for item in ['当月','当年']:
        shapes = getshape(title_left,item+'新续转增速')
        df = pd.DataFrame(cursor.fetchall())
        df_values = pct_2_float(div_10000(df)).iloc[:-1,:]
        for xbflag, left_inches in {'新保':0.2, '续保':5.0, '转保':10}.items():
            df_bar = df_values[['部门', xbflag+'增速']].sort_values(by=xbflag+'增速')
            drawbar(shapes, df_bar, item + xbflag + '增速', left=Inches(left_inches), top=Inches(1), width=Inches(4), height=Inches(6.5),
                    fontsize=14)
    # 当月小口径续保率、转保率
    # shapes = getshape(title_left, item + '新续转增速')
    df = pd.DataFrame(cursor.fetchall())
    slide_with_bar_and_table(df, '当月各部门续保率', '当月各部门小口径续保率', ['部门', '当月小口径续保率', '当月小口径同比', '当月大口径续保率','当月大口径同比'], ['部门', '当月小口径续保率'],
                             title_left, '当月续保率')
    slide_with_bar_and_table(df, '当月各部门转保率', '当月各部门转保率', ['部门', '当月转保率', '当月转保率同比'], ['部门', '当月转保率'],
                             title_left, '当月转保率')
    slide_with_bar_and_table(df, '当年各部门续保率', '当年各部门小口径续保率', ['部门', '当年小口径续保率', '当年小口径同比', '当年大口径续保率', '当年大口径同比'],
                             ['部门', '当年小口径续保率'],
                             title_left, '当年续保率')
    slide_with_bar_and_table(df, '当年各部门转保率', '当年各部门转保率', ['部门', '当年转保率', '当年转保率同比'], ['部门', '当年转保率'],
                             title_left, '当年转保率')
    df = pd.DataFrame(cursor.fetchall())
    slide_with_bar_and_table(df, '当月各部门家用车续保率', '当月各部门家用车小口径续保率', ['部门', '当月小口径续保率', '当月小口径同比', '当月大口径续保率', '当月大口径同比'],
                             ['部门', '当月小口径续保率'],
                             title_left, '当月续保率--家用车')
    slide_with_bar_and_table(df, '当月各部门家用车转保率', '当月各部门家用车转保率', ['部门', '当月转保率', '当月转保率同比'], ['部门', '当月转保率'],
                             title_left, '当月转保率--家用车')
    slide_with_bar_and_table(df, '当年各部门家用车续保率', '当年各部门家用车小口径续保率', ['部门', '当年小口径续保率', '当年小口径同比', '当年大口径续保率', '当年大口径同比'],
                             ['部门', '当年小口径续保率'],
                             title_left, '当年续保率--家用车')
    slide_with_bar_and_table(df, '当年各部门家用车转保率', '当年各部门家用车转保率', ['部门', '当年转保率', '当年转保率同比'], ['部门', '当年转保率'],
                             title_left, '当年转保率--家用车')
    # 提前签单跳过
    df = pd.DataFrame(cursor.fetchall())
    # 当月/当年渠道增速
    for item in ['当月','当年']:
        shapes = getshape(title_left,item+'各部门分渠道增速')
        df = pd.DataFrame(cursor.fetchall())
        df_values = df.replace('/','').applymap(format_pct).sort_index(axis=1,ascending=False)
        drawtable(shapes,df_values,item+'各部门分渠道增速',Inches(0.5),Inches(1.5),Inches(13))
    # 标题页3
    getshape('03', '近期重点工作', 3)
    outfn = year_str+month_str+'车险业务汇报--珠海.pptx'
    prs.save(outfn)


orange = RGBColor(210, 71, 38)  # 橙色
black = RGBColor(0, 0, 0)  # 黑色
red = RGBColor(153, 0, 0)  # 砖红色
templatename = '2018template.pptx'  # 模板名称
prs = Presentation(templatename)
delta = datetime.timedelta(days=1)
today = datetime.date.today()
datelist = get_datelist(today)  # 获取1月到上一个完整月的最后一天
deskdate_date = dateutil.parser.parse(datelist[-1])  #　上一个完整月的最后一天，datetime格式
# deskdate_date = today-delta
# deskdate_date = datetime.date(2018,1,5)
deskdate = datetime.date.strftime((deskdate_date+delta),'%Y%m%d')  # 字符串格式
# deskdate = '20190201'
year_str = str(deskdate_date.year)+'年'
month_str = str(deskdate_date.month)+'月'
allmonth_str = '1-%s' % month_str if month_str != '1月' else month_str  # 累计月份，例：1-5月
user='isp'
password='pIcc4404'
host='56.49.33.246:1433'
db='CAR'
engine=create_engine(str(r"mssql+pymssql://%s:"+'%s'+"@%s/%s")%(user,password,host,db))

main()



